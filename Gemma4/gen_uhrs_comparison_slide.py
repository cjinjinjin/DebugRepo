"""Generate PPT: UHRS Human Label Comparison — Gemma4 Two-Step vs One-Step vs GPT5."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------- colors ----------
DARK_BG        = RGBColor(0x1E, 0x1E, 0x2E)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xB0, 0xB0, 0xB0)
ACCENT_BLUE    = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_GREEN   = RGBColor(0x66, 0xBB, 0x6A)
ACCENT_ORANGE  = RGBColor(0xFF, 0xA7, 0x26)
ACCENT_RED     = RGBColor(0xEF, 0x53, 0x50)
TABLE_HEADER_BG = RGBColor(0x33, 0x33, 0x55)
TABLE_ROW_BG    = RGBColor(0x2A, 0x2A, 0x40)
TABLE_ALT_BG    = RGBColor(0x25, 0x25, 0x3A)
HIGHLIGHT_BG    = RGBColor(0x1B, 0x3A, 0x2A)


def set_cell(cell, text, size=9, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)


def fill_cell_bg(cell, color):
    from lxml import etree
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgbClr.set('val', '%02x%02x%02x' % (color[0] if isinstance(color, tuple) else color.red,
                                           color[1] if isinstance(color, tuple) else color.green,
                                           color[2] if isinstance(color, tuple) else color.blue))


def add_textbox(slide, left, top, width, height, text, size=11, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_rich_textbox(slide, left, top, width, height, lines, default_size=10):
    """lines: list of (text, size, bold, color)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size or default_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


def make_table(slide, left, top, width, height, headers, data_rows, header_colors=None):
    """Create a styled table. header_colors: list of RGBColor per column header."""
    nrows = 1 + len(data_rows)
    ncols = len(headers)
    tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table

    for c, h in enumerate(headers):
        hc = header_colors[c] if header_colors else ACCENT_BLUE
        set_cell(tbl.cell(0, c), h, size=9, bold=True, color=hc)
        fill_cell_bg(tbl.cell(0, c), TABLE_HEADER_BG)

    for r, row in enumerate(data_rows):
        bg = TABLE_ROW_BG if r % 2 == 0 else TABLE_ALT_BG
        for c, (val, color, bold) in enumerate(row):
            set_cell(tbl.cell(r + 1, c), val, size=9, bold=bold, color=color)
            fill_cell_bg(tbl.cell(r + 1, c), bg)
    return tbl


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # ===== TITLE =====
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
                "UHRS Human Label Quality Comparison",
                size=24, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.3),
                "Gemma4 Two-Step  vs  Gemma4 One-Step  vs  GPT5  |  Random 200 LPs  |  3 Judges per Image (Max-Vote)",
                size=11, color=LIGHT_GRAY)

    # ===== TABLE 1: Image Level Good Rate =====
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(5.5), Inches(0.3),
                "1. Image Level Quality (Max-Vote)", size=13, bold=True, color=ACCENT_BLUE)

    W = WHITE
    G = ACCENT_GREEN
    O = ACCENT_ORANGE
    R = ACCENT_RED

    headers1 = ["Metric", "Gemma4\nTwo-Step", "Gemma4\nOne-Step", "GPT5", "Two-Step\nvs GPT5"]
    data1 = [
        [("Total Images", W, False), ("994",        W, False), ("994",        W, False), ("998",       W, False), ("—",         W, False)],
        [("Good",         W, True),  ("755 (76.0%)", G, True),  ("749 (75.3%)", G, False), ("653 (65.4%)", W, False), ("+10.6pp",   G, True)],
        [("Fair",         W, False), ("94 (9.5%)",   W, False), ("73 (7.3%)",   W, False), ("65 (6.5%)",  W, False), ("+3.0pp",    W, False)],
        [("Bad",          W, False), ("145 (14.6%)", W, False), ("172 (17.3%)", O, False), ("280 (28.1%)", R, False), ("-13.5pp",   G, True)],
    ]
    make_table(slide, Inches(0.5), Inches(1.45), Inches(6.2), Inches(1.5), headers1, data1)

    # ===== TABLE 2: LP Level N/5 Good Distribution =====
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(6), Inches(0.3),
                "2. LP Level — N/5 Good Distribution", size=13, bold=True, color=ACCENT_BLUE)

    headers2 = ["N/5 Good", "Gemma4\nTwo-Step", "Gemma4\nOne-Step", "GPT5"]
    data2 = [
        [("0/5", W, False), ("1 (0.5%)",   W, False), ("2 (1.0%)",   W, False), ("1 (0.5%)",   W, False)],
        [("1/5", W, False), ("4 (2.0%)",   W, False), ("5 (2.5%)",   W, False), ("9 (4.5%)",   O, False)],
        [("2/5", W, False), ("21 (10.5%)", W, False), ("18 (9.0%)",  W, False), ("39 (19.5%)", O, False)],
        [("3/5", W, False), ("46 (23.0%)", W, False), ("49 (24.5%)", W, False), ("65 (32.5%)", W, False)],
        [("4/5", W, False), ("69 (34.5%)", G, False), ("69 (34.5%)", G, False), ("59 (29.5%)", W, False)],
        [("5/5", W, False), ("59 (29.5%)", G, True),  ("57 (28.5%)", G, False), ("27 (13.5%)", W, False)],
    ]
    make_table(slide, Inches(0.5), Inches(3.45), Inches(5.5), Inches(2.2), headers2, data2)

    # ===== TABLE 3: LP Level Cumulative =====
    add_textbox(slide, Inches(6.8), Inches(1.1), Inches(6), Inches(0.3),
                "3. LP Level — Cumulative Good Rate", size=13, bold=True, color=ACCENT_BLUE)

    headers3 = ["Threshold", "Gemma4\nTwo-Step", "Gemma4\nOne-Step", "GPT5", "Two-Step\nvs GPT5"]
    data3 = [
        [("≥ 1/5", W, False), ("99.5%", W, False), ("99.0%", W, False), ("99.5%", W, False), ("+0.0pp",  W, False)],
        [("≥ 2/5", W, False), ("97.0%", W, False), ("96.5%", W, False), ("95.0%", W, False), ("+2.0pp",  W, False)],
        [("≥ 3/5", W, False), ("87.0%", G, False), ("87.5%", G, False), ("75.5%", W, False), ("+11.5pp", G, True)],
        [("≥ 4/5", W, False), ("64.0%", G, True),  ("63.0%", G, False), ("43.0%", W, False), ("+21.0pp", G, True)],
        [("≥ 5/5", W, False), ("29.5%", G, True),  ("28.5%", G, False), ("13.5%", W, False), ("+16.0pp", G, True)],
    ]
    make_table(slide, Inches(6.8), Inches(1.45), Inches(6.0), Inches(1.7), headers3, data3)

    # ===== TABLE 4: Speed & Prompt Comparison =====
    add_textbox(slide, Inches(6.8), Inches(3.3), Inches(6), Inches(0.3),
                "4. Speed & Prompt Length", size=13, bold=True, color=ACCENT_BLUE)

    headers4 = ["Metric", "Gemma4\nTwo-Step", "Gemma4\nOne-Step", "GPT5"]
    data4 = [
        [("Avg Latency",       W, False), ("84.9s",     G, True),  ("122.7s",    O, False), ("—",         LIGHT_GRAY, False)],
        [("Avg Prompt Words",  W, False), ("42 words",  G, True),  ("127 words", O, False), ("—",         LIGHT_GRAY, False)],
        [("Speed Gain",        W, False), ("-30.8%",    G, True),  ("baseline",  W, False), ("—",         LIGHT_GRAY, False)],
    ]
    make_table(slide, Inches(6.8), Inches(3.65), Inches(5.5), Inches(1.2), headers4, data4)

    # ===== KEY TAKEAWAYS =====
    add_textbox(slide, Inches(6.8), Inches(5.05), Inches(6), Inches(0.3),
                "5. Key Takeaways", size=13, bold=True, color=ACCENT_ORANGE)

    findings = [
        ("●  Gemma4 >> GPT5 across all metrics", 10, True, ACCENT_GREEN),
        ("     Image Good Rate: 76.0% vs 65.4% (+10.6pp)", 9, False, WHITE),
        ("     LP ≥4/5 Good: 64.0% vs 43.0% (+21.0pp)", 9, False, WHITE),
        ("", 6, False, WHITE),
        ("●  Two-Step ≈ One-Step in quality", 10, True, ACCENT_BLUE),
        ("     Image Good Rate: 76.0% vs 75.3% (+0.7pp)", 9, False, WHITE),
        ("     Bad Rate lower: 14.6% vs 17.3% (-2.7pp)", 9, False, WHITE),
        ("", 6, False, WHITE),
        ("●  Two-Step wins on speed & prompt control", 10, True, ACCENT_ORANGE),
        ("     84.9s vs 122.7s (-30.8%), 42 vs 127 words (-67%)", 9, False, WHITE),
    ]
    add_rich_textbox(slide, Inches(6.8), Inches(5.35), Inches(6), Inches(2.0), findings)

    # ===== BOTTOM RECOMMENDATION =====
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(5.8), Inches(1.2),
                "Recommendation\n"
                "Two-Step is the preferred method:\n"
                "  • Same quality as One-Step\n"
                "  • 31% faster inference\n"
                "  • 67% shorter prompts (better T2I quality)\n"
                "  • Both significantly outperform GPT5",
                size=10, bold=False, color=WHITE)

    # Highlight box for recommendation title
    add_textbox(slide, Inches(0.5), Inches(5.7), Inches(2), Inches(0.3),
                "✦ Recommendation", size=12, bold=True, color=ACCENT_ORANGE)

    out_path = "Gemma4/UHRS_Quality_Comparison_ThreeModels.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
