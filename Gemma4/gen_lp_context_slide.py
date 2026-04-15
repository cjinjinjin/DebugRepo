"""Generate a single-slide PPT: LP Context Distribution & Speed Impact."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ---------- data ----------

# LP content length distribution (sft_eval_cot, 87 samples)
length_dist = {
    "Dataset": ["sft_eval_cot", "sft_train_cot", "grpo_train"],
    "Samples": [87, 833, 1100],
    "Avg": [4689, 5186, 3867],
    "Median": [4055, 3779, 2890],
    "Max": [17187, 343978, 35687],
}

# Truncation coverage (sft_eval_cot)
trunc_buckets = [500, 1000, 2000, 3000, 5000, 10000]
trunc_pct = [1.1, 8.0, 16.1, 32.2, 63.2, 92.0]

# LP chars sweep → speed (single GPU BF16 No-CoT, 20 samples)
lp_sweep = {
    "LP Chars":   ["400", "1,000", "2,000", "5,000", "Unlimited"],
    "Avg Input Tok": [533, 636, 779, 1032, 1095],
    "Avg TTFT (s)":  [0.02, 0.02, 0.02, 0.02, 0.02],
    "Avg Decode (s)": [54.8, 53.9, 54.1, 51.7, 52.2],
    "Avg Total (s)":  [54.8, 53.9, 54.1, 51.7, 52.2],
    "tok/s":          [11.8, 11.7, 11.7, 11.9, 12.2],
}

# ---------- colors ----------
DARK_BG   = RGBColor(0x1E, 0x1E, 0x2E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB0)
ACCENT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
ACCENT_ORANGE = RGBColor(0xFF, 0xA7, 0x26)
TABLE_HEADER_BG = RGBColor(0x33, 0x33, 0x55)
TABLE_ROW_BG    = RGBColor(0x2A, 0x2A, 0x40)
TABLE_ALT_BG    = RGBColor(0x25, 0x25, 0x3A)


def set_cell(cell, text, size=9, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Remove cell margins for compactness
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)


def fill_cell_bg(cell, color):
    from lxml import etree
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgbClr.set('val', '%02x%02x%02x' % (color[0] if isinstance(color, tuple) else color.red,
                                           color[1] if isinstance(color, tuple) else color.green,
                                           color[2] if isinstance(color, tuple) else color.blue))


def add_textbox(slide, left, top, width, height, text, size=11, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # ===== TITLE =====
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                "LP Context Distribution & Impact on Inference Speed",
                size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.35),
                "Gemma 4 26B-A4B-it  |  BF16 No-CoT  |  Single A100-SXM4-80GB",
                size=12, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # ===== LEFT SECTION: LP Content Length Distribution =====
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.35),
                "1. LP Content Length Distribution (chars)",
                size=14, bold=True, color=ACCENT_BLUE)

    # Table 1: Length stats
    rows, cols = 4, 6
    tbl1 = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(5.8), Inches(1.2)).table
    headers1 = ["Dataset", "Samples", "Avg", "Median", "P95", "Max"]
    for c, h in enumerate(headers1):
        set_cell(tbl1.cell(0, c), h, size=9, bold=True, color=ACCENT_BLUE)
        fill_cell_bg(tbl1.cell(0, c), TABLE_HEADER_BG)

    data_rows = [
        ["sft_eval_cot", "87", "4,689", "4,055", "10,988", "17,187"],
        ["sft_train_cot", "833", "5,186", "3,779", "—", "343,978"],
        ["grpo_train", "1,100", "3,867", "2,890", "—", "35,687"],
    ]
    for r, row_data in enumerate(data_rows):
        bg_color = TABLE_ROW_BG if r % 2 == 0 else TABLE_ALT_BG
        for c, val in enumerate(row_data):
            set_cell(tbl1.cell(r+1, c), val, size=9)
            fill_cell_bg(tbl1.cell(r+1, c), bg_color)

    # Table 2: Truncation coverage
    add_textbox(slide, Inches(0.5), Inches(2.9), Inches(5.5), Inches(0.3),
                "Truncation Coverage (sft_eval_cot, 87 samples)",
                size=11, bold=True, color=ACCENT_GREEN)

    rows2, cols2 = 3, 6
    tbl2 = slide.shapes.add_table(rows2, cols2, Inches(0.5), Inches(3.25), Inches(5.8), Inches(0.85)).table
    thresholds = ["≤500", "≤1,000", "≤2,000", "≤3,000", "≤5,000", "≤10,000"]
    pcts = ["1.1%", "8.0%", "16.1%", "32.2%", "63.2%", "92.0%"]
    for c, h in enumerate(thresholds):
        set_cell(tbl2.cell(0, c), h, size=9, bold=True, color=ACCENT_GREEN)
        fill_cell_bg(tbl2.cell(0, c), TABLE_HEADER_BG)
    for c, h in enumerate(["Chars ≤", "Chars ≤", "Chars ≤", "Chars ≤", "Chars ≤", "Chars ≤"]):
        pass  # header row already set
    # Label row
    for c in range(cols2):
        set_cell(tbl2.cell(1, c), "Coverage", size=8, color=LIGHT_GRAY)
        fill_cell_bg(tbl2.cell(1, c), TABLE_ROW_BG)
    for c, val in enumerate(pcts):
        set_cell(tbl2.cell(2, c), val, size=10, bold=True,
                 color=ACCENT_ORANGE if c <= 2 else ACCENT_GREEN)
        fill_cell_bg(tbl2.cell(2, c), TABLE_ALT_BG)

    # Insight callout
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(5.8), Inches(0.6),
                "84% of samples exceed 2,000 chars. Median ~4K chars.\n"
                "Core product info typically within first 1,000–2,000 chars.",
                size=10, bold=False, color=LIGHT_GRAY)

    # ===== RIGHT SECTION: Speed Impact =====
    add_textbox(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(0.35),
                "2. LP Truncation vs Inference Speed",
                size=14, bold=True, color=ACCENT_BLUE)

    # Table 3: Speed sweep
    rows3, cols3 = 6, 6
    tbl3 = slide.shapes.add_table(rows3, cols3, Inches(6.8), Inches(1.6), Inches(6.0), Inches(1.8)).table
    headers3 = ["LP Chars", "Avg Input\nTokens", "TTFT (s)", "Decode (s)", "Total (s)", "tok/s"]
    for c, h in enumerate(headers3):
        set_cell(tbl3.cell(0, c), h, size=9, bold=True, color=ACCENT_BLUE)
        fill_cell_bg(tbl3.cell(0, c), TABLE_HEADER_BG)

    speed_data = [
        ["400",       "533",  "0.02", "54.8", "54.8", "11.8"],
        ["1,000",     "636",  "0.02", "53.9", "53.9", "11.7"],
        ["2,000",     "779",  "0.02", "54.1", "54.1", "11.7"],
        ["5,000",     "1,032","0.02", "51.7", "51.7", "11.9"],
        ["Unlimited", "1,095","0.02", "52.2", "52.2", "12.2"],
    ]
    for r, row_data in enumerate(speed_data):
        bg_color = TABLE_ROW_BG if r % 2 == 0 else TABLE_ALT_BG
        for c, val in enumerate(row_data):
            color = WHITE
            if r == 4:  # highlight unlimited row
                color = ACCENT_GREEN
            set_cell(tbl3.cell(r+1, c), val, size=9, bold=(r == 4), color=color)
            fill_cell_bg(tbl3.cell(r+1, c), bg_color)

    # ===== Key Findings =====
    add_textbox(slide, Inches(6.8), Inches(3.6), Inches(6), Inches(0.35),
                "3. Key Findings",
                size=14, bold=True, color=ACCENT_ORANGE)

    findings = (
        "●  Input tokens 2x (533→1,095) but total time flat (~53s)\n"
        "     → Bottleneck is 100% decode, NOT prefill\n\n"
        "●  TTFT constant at 0.02s regardless of input length\n"
        "     → Prefill cost negligible on single A100\n\n"
        "●  tok/s stable at 11.7–12.2 across all LP lengths\n"
        "     → LP truncation does NOT improve throughput\n\n"
        "●  Truncation value: defensive only for extreme tail\n"
        "     (P95=10,988 chars, max=343K chars in train set)\n"
        "     Prevents OOM / KV-cache overflow on outliers"
    )
    add_textbox(slide, Inches(6.8), Inches(4.0), Inches(6), Inches(3.0),
                findings, size=10, bold=False, color=WHITE)

    # ===== Bottom bar =====
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.35),
                "Recommendation:  Keep --max_lp_chars 2000 for safety (OOM defense on tail), "
                "but expect no speed gain on average cases.  Focus optimization on decode (quantization / shorter output).",
                size=10, bold=True, color=ACCENT_ORANGE, align=PP_ALIGN.LEFT)

    # Save
    out_path = "Gemma4/LP_Context_Distribution_Speed_Impact.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
